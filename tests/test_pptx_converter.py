"""PPTX → PDF 轉換測試。

無 PowerPoint 也無 LibreOffice 的環境整個 skip。
CI / 乾淨機器上此測試會顯示為 skipped，不會失敗。
"""

from __future__ import annotations

import shutil
import sys
from pathlib import Path

import pytest


def _any_converter_available() -> bool:
    if sys.platform == "win32":
        try:
            import comtypes.client  # type: ignore  # noqa: F401
            # 有 comtypes 不代表裝了 PowerPoint；假設有
            return True
        except ImportError:
            pass
    if shutil.which("soffice") or shutil.which("soffice.exe"):
        return True
    # Windows 常見路徑
    if sys.platform == "win32":
        for base in ("C:/Program Files/LibreOffice/program/soffice.exe",
                     "C:/Program Files (x86)/LibreOffice/program/soffice.exe"):
            if Path(base).exists():
                return True
    return False


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """用 python-pptx 產生最小 PPTX。若環境無 python-pptx 則跳過。"""
    pptx_lib = pytest.importorskip("pptx")
    from pptx.util import Inches

    prs = pptx_lib.Presentation()
    for title in ("Slide A", "Slide B"):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
    out = tmp_path / "sample.pptx"
    prs.save(str(out))
    return out


def test_cache_key_stable_for_same_file(tmp_path):
    """同一檔案路徑產 key 應穩定。"""
    from teleprompter.core.pptx_converter import _cache_key

    f = tmp_path / "a.pptx"
    f.write_bytes(b"x" * 100)
    k1 = _cache_key(f)
    k2 = _cache_key(f)
    assert k1 == k2
    # 改變內容 + mtime → key 改變
    import time
    time.sleep(0.05)
    f.write_bytes(b"x" * 200)
    assert _cache_key(f) != k1


def test_find_libreoffice_returns_path_or_none():
    """純 smoke 測試：回傳值是 str 或 None，不丟例外。"""
    from teleprompter.core.pptx_converter import _find_libreoffice

    result = _find_libreoffice()
    assert result is None or isinstance(result, str)


@pytest.mark.skipif(not _any_converter_available(),
                    reason="需要 PowerPoint 或 LibreOffice 才能測")
def test_convert_pptx_to_pdf_produces_file(sample_pptx):
    from teleprompter.core.pptx_converter import convert_pptx_to_pdf

    pdf = convert_pptx_to_pdf(sample_pptx)
    assert pdf.exists()
    assert pdf.suffix == ".pdf"
    assert pdf.stat().st_size > 0


def test_convert_missing_file_raises(tmp_path):
    from teleprompter.core.pptx_converter import convert_pptx_to_pdf

    with pytest.raises(FileNotFoundError):
        convert_pptx_to_pdf(tmp_path / "nope.pptx")


def test_convert_wrong_extension_raises(tmp_path):
    from teleprompter.core.pptx_converter import convert_pptx_to_pdf

    fake = tmp_path / "a.txt"
    fake.write_text("not pptx")
    with pytest.raises(ValueError):
        convert_pptx_to_pdf(fake)

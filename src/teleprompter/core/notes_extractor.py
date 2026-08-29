"""從投影片檔自動抽取講稿（speaker notes）。

支援來源：

1. **PPTX**：讀每頁的備忘稿（`slide.notes_slide`）— 最準確，作者本來就寫在那裡
2. **HTML**（Claude Design / reveal.js 等）：讀 `data-speaker-notes` 屬性，
   或 `<aside class="notes">` 區塊
3. **PDF**：沒有 notes 概念，退而取每頁的可見文字當草稿

抽出的結果會組成標準講稿格式（`# Slide N` 標題 + 內文 + `---` 分頁），
可直接餵給 `transcript_loader.load_transcript()`。

設計原則：抽不到就回空字串那一頁，不臆造內容——讓使用者在預覽視窗自行補。
"""

from __future__ import annotations

import html as html_mod
import re
from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class ExtractedNotes:
    """抽取結果。"""

    source: str = ""                       # 來源類型：pptx / html / pdf
    pages: list[str] = field(default_factory=list)   # 每頁講稿（可能為空字串）
    titles: list[str] = field(default_factory=list)  # 每頁標題（沒有就空字串）

    @property
    def page_count(self) -> int:
        return len(self.pages)

    @property
    def filled_count(self) -> int:
        """實際抽到內容的頁數。"""
        return sum(1 for p in self.pages if p.strip())

    def to_transcript_text(self) -> str:
        """組成標準講稿格式：`# Slide N · 標題` + 內文，`---` 分頁。"""
        blocks = []
        for i, body in enumerate(self.pages, 1):
            title = self.titles[i - 1] if i - 1 < len(self.titles) else ""
            heading = f"# Slide {i}" + (f" · {title}" if title else "")
            blocks.append(heading + "\n\n" + (body.strip() or ""))
        return "\n\n---\n\n".join(blocks).strip() + "\n"


def _clean(text: str) -> str:
    """去掉多餘空白，保留段落換行。"""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    lines = [ln.rstrip() for ln in text.split("\n")]
    out: list[str] = []
    for ln in lines:
        if ln.strip() or (out and out[-1].strip()):
            out.append(ln.strip())
    return "\n".join(out).strip()


# ============================================================
# PPTX：讀 notes_slide
# ============================================================

def extract_from_pptx(path: str | Path) -> ExtractedNotes:
    from pptx import Presentation

    prs = Presentation(str(path))
    result = ExtractedNotes(source="pptx")
    for slide in prs.slides:
        note = ""
        if slide.has_notes_slide:
            frame = slide.notes_slide.notes_text_frame
            if frame is not None and frame.text:
                note = _clean(frame.text)
        result.pages.append(note)
        result.titles.append(_pptx_slide_title(slide))
    return result


def _pptx_slide_title(slide) -> str:
    try:
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            return _clean(slide.shapes.title.text).split("\n")[0][:60]
    except Exception:
        pass
    # 沒有標題版面 → 取第一個有字的文字框
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
            return _clean(shape.text_frame.text).split("\n")[0][:60]
    return ""


# ============================================================
# HTML：讀 data-speaker-notes / <aside class="notes">
# ============================================================

_SECTION_RE = re.compile(r"<section\b[^>]*>.*?</section>", re.S | re.I)
_NOTES_ATTR_RE = re.compile(r'data-speaker-notes="([^"]*)"', re.I)
_LABEL_ATTR_RE = re.compile(r'data-label="([^"]*)"', re.I)
_ASIDE_RE = re.compile(r'<aside[^>]*class="[^"]*notes[^"]*"[^>]*>(.*?)</aside>', re.S | re.I)
_TAG_RE = re.compile(r"<[^>]+>")
# 講稿裡常見的私有註記段落（來源標註），投影時不需要唸出來
_PRIVATE_SECTION_RE = re.compile(r"\n?\s*\[(?:Sources|Stage|Notes)\][\s\S]*$", re.I)


def extract_from_html(path: str | Path) -> ExtractedNotes:
    raw = Path(path).read_text(encoding="utf-8", errors="ignore")
    result = ExtractedNotes(source="html")
    for sec in _SECTION_RE.findall(raw):
        note = ""
        m = _NOTES_ATTR_RE.search(sec)
        if m:
            note = html_mod.unescape(m.group(1))
        else:
            a = _ASIDE_RE.search(sec)
            if a:
                note = html_mod.unescape(_TAG_RE.sub(" ", a.group(1)))
        note = _PRIVATE_SECTION_RE.sub("", note)
        result.pages.append(_clean(note))
        lab = _LABEL_ATTR_RE.search(sec)
        result.titles.append(html_mod.unescape(lab.group(1)).strip() if lab else "")
    return result


# ============================================================
# PDF：沒有 notes → 取頁面可見文字當草稿
# ============================================================

def extract_from_pdf(path: str | Path) -> ExtractedNotes:
    import fitz

    result = ExtractedNotes(source="pdf")
    doc = fitz.open(str(path))
    try:
        for page in doc:
            text = _clean(page.get_text() or "")
            lines = [ln for ln in text.split("\n") if ln.strip()]
            title = lines[0][:60] if lines else ""
            body = "\n".join(lines[1:]) if len(lines) > 1 else ""
            result.pages.append(body)
            result.titles.append(title)
    finally:
        doc.close()
    return result


# ============================================================
# 主入口
# ============================================================

def extract_notes(path: str | Path) -> ExtractedNotes:
    """依副檔名選擇抽取器；不支援的格式回空結果。"""
    p = Path(path)
    suffix = p.suffix.lower()
    if suffix in (".pptx", ".ppt"):
        return extract_from_pptx(p)
    if suffix in (".html", ".htm"):
        return extract_from_html(p)
    if suffix == ".pdf":
        return extract_from_pdf(p)
    return ExtractedNotes()
